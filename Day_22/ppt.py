from email.quoprimime import body_check

from numpy.ma.core import shape
from  pptx import Presentation
pres = Presentation()

title_slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(title_slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Welcome to Python"
subtitle.text = "life is short,you must to be dead"

bullet_slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(bullet_slide_layout)

shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Introduction"

tf = body_shape.text_frame
tf.text = "History of Python"

p = tf.add_paragraph()
p.text = 'X\'max 1989'
p.level = 3
# 添加一个二级段落
p = tf.add_paragraph()
p.text = 'Guido began to write interpreter for Python.'
p.level = 2

# 保存幻灯片
pres.save('test.pptx')